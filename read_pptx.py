import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

pptx_path = r"C:\Users\Admin\Downloads\ITSS.pptx"
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*60}")
    print(f"SLIDE {i}")
    print(f"{'='*60}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(" | ".join(cells))
